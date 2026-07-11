"""
Antiparallel Flow — Strategy Deck Generator
Generates a board/investor-ready pitch deck (.pptx)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- Brand tokens (from Antiparallel design system) ----------
BG_DARK = RGBColor(0x0A, 0x0A, 0x0A)
BG_SURFACE = RGBColor(0x14, 0x14, 0x14)
BG_SURFACE_2 = RGBColor(0x1F, 0x1F, 0x1F)
TEXT_PRIMARY = RGBColor(0xFA, 0xFA, 0xFA)
TEXT_SECONDARY = RGBColor(0xA1, 0xA1, 0xA1)
TEXT_MUTED = RGBColor(0x71, 0x71, 0x71)
MINT = RGBColor(0x00, 0xFF, 0xBA)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
RED = RGBColor(0xFF, 0x47, 0x57)
BORDER = RGBColor(0x2A, 0x2A, 0x2A)

FONT_HEAD = "Inter"
FONT_BODY = "Inter"
FONT_THAI = "IBM Plex Sans Thai"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_fill(shp, color)
    if not line:
        shp.line.fill.background()
    return shp


def add_text(slide, x, y, w, h, text, size=18, color=TEXT_PRIMARY,
             bold=False, font=FONT_HEAD, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bg(slide, color=BG_DARK):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)


def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(4), Inches(0.3),
             "Antiparallel Flow  ·  Strategy Deck  ·  Confidential",
             size=9, color=TEXT_MUTED, font=FONT_HEAD)
    add_text(slide, Inches(12.0), Inches(7.1), Inches(0.8), Inches(0.3),
             f"{page_num:02d} / {total:02d}", size=9, color=TEXT_MUTED,
             font=FONT_HEAD, align=PP_ALIGN.RIGHT)


def add_section_chip(slide, x, y, label, accent=MINT):
    chip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.18), Inches(0.18))
    set_fill(chip, accent)
    add_text(slide, x + Inches(0.32), y - Inches(0.02), Inches(4), Inches(0.3),
             label.upper(), size=10, color=accent, bold=True, font=FONT_HEAD)


def add_divider(slide, x, y, w, color=BORDER):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(9525))
    set_fill(line, color)


# ---------- Slide builders ----------

def slide_blank(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def title_slide(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    # accent bar
    add_rect(s, 0, 0, Inches(0.15), SLIDE_H, MINT)
    # eyebrow
    add_text(s, Inches(0.8), Inches(0.8), Inches(6), Inches(0.4),
             "ANTIPARALLEL  ·  INTERNAL STRATEGY",
             size=11, color=MINT, bold=True, font=FONT_HEAD)
    # logo block
    add_text(s, Inches(0.8), Inches(2.6), Inches(12), Inches(1.4),
             "Antiparallel Flow",
             size=72, color=TEXT_PRIMARY, bold=True, font=FONT_HEAD)
    add_text(s, Inches(0.8), Inches(4.1), Inches(12), Inches(0.6),
             "AI Workflow Automation built for Thai businesses",
             size=24, color=TEXT_SECONDARY, font=FONT_HEAD)
    add_text(s, Inches(0.8), Inches(4.8), Inches(12), Inches(0.5),
             "เปลี่ยนความรู้ AI ในห้องเรียน ให้กลายเป็นงานจริงในองค์กร",
             size=18, color=MINT, font=FONT_THAI)
    # bottom meta
    add_divider(s, Inches(0.8), Inches(6.6), Inches(11.7))
    add_text(s, Inches(0.8), Inches(6.75), Inches(6), Inches(0.3),
             "Board / Investor Strategy Deck", size=11, color=TEXT_SECONDARY, font=FONT_HEAD)
    add_text(s, Inches(8.5), Inches(6.75), Inches(4), Inches(0.3),
             "May 2026  ·  v1.0", size=11, color=TEXT_SECONDARY,
             font=FONT_HEAD, align=PP_ALIGN.RIGHT)


def section_divider(prs, section_num, title_en, title_th, page, total, accent=MINT):
    s = slide_blank(prs)
    add_bg(s, BG_SURFACE)
    add_rect(s, 0, 0, Inches(0.15), SLIDE_H, accent)
    add_text(s, Inches(0.8), Inches(2.8), Inches(3), Inches(0.6),
             f"SECTION {section_num:02d}", size=14, color=accent, bold=True, font=FONT_HEAD)
    add_text(s, Inches(0.8), Inches(3.4), Inches(12), Inches(1.0),
             title_en, size=56, color=TEXT_PRIMARY, bold=True, font=FONT_HEAD)
    add_text(s, Inches(0.8), Inches(4.6), Inches(12), Inches(0.6),
             title_th, size=22, color=TEXT_SECONDARY, font=FONT_THAI)
    add_footer(s, page, total)


def content_header(slide, eyebrow, title, accent=MINT):
    add_section_chip(slide, Inches(0.6), Inches(0.55), eyebrow, accent)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
             title, size=32, color=TEXT_PRIMARY, bold=True, font=FONT_HEAD)
    add_divider(slide, Inches(0.6), Inches(1.85), Inches(12.1))


def slide_executive_summary(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Executive Summary", "What this deck argues, in one page")

    points = [
        ("01", "PROBLEM",
         "Corporates spend on AI training, but knowledge dies the moment\nemployees return to a workflow that doesn't support AI."),
        ("02", "INSIGHT",
         "Global AI tools (ChatGPT/Claude) can't reach Thai systems —\nLine OA, FlowAccount, PromptPay, สรรพากร. We can."),
        ("03", "PRODUCT",
         "Antiparallel Flow — AI workflow automation that triggers\nfrom Thai-specific systems and runs real work, not chat."),
        ("04", "GTM",
         "Distribute through existing LMS corporate customers.\nDone-for-you in months 1–2, then productize into Workflow Packs."),
        ("05", "ECONOMICS",
         "ARR per deal ฿420k (mid) to ฿2.4M (enterprise).\nTarget ฿7–10M ARR by Month 6, 20+ paying customers."),
        ("06", "MOAT",
         "Thai integrations + workflow library + LMS distribution.\nModel-provider-proof: smarter models = better product."),
    ]
    col_w = Inches(4.0)
    row_h = Inches(2.3)
    start_x = Inches(0.6)
    start_y = Inches(2.15)
    gap = Inches(0.15)
    for i, (num, label, body) in enumerate(points):
        col = i % 3
        row = i // 3
        x = start_x + (col_w + gap) * col
        y = start_y + (row_h + gap) * row
        card = add_rect(s, x, y, col_w, row_h, BG_SURFACE)
        add_text(s, x + Inches(0.3), y + Inches(0.25), Inches(1.5), Inches(0.4),
                 num, size=12, color=MINT, bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.3), y + Inches(0.6), Inches(3.5), Inches(0.4),
                 label, size=14, color=TEXT_PRIMARY, bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.3), y + Inches(1.05), Inches(3.5), Inches(1.2),
                 body, size=11, color=TEXT_SECONDARY, font=FONT_HEAD)
    add_footer(s, page, total)


def slide_tldr(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "TL;DR", "ถ้าอ่านได้แค่หน้าเดียว อ่านหน้านี้")
    body = [
        ("เราเป็นใคร", "Antiparallel — LMS สอน AI ให้องค์กรไทย มีลูกค้า corporate active อยู่แล้ว"),
        ("ปัญหา", "ลูกค้าจ่ายค่าเทรน AI หลักแสน-ล้าน แต่พนักงานกลับไปทำงานเดิม ไม่มีระบบให้ใช้ AI ต่อ"),
        ("Insight", "ChatGPT/Claude ตามไม่ทันระบบไทย (Line OA, FlowAccount, ใบกำกับภาษี) — ช่องว่างนี้ใครเข้าก็ผูกขาด"),
        ("ทางออก", "Antiparallel Flow — workflow automation ที่ trigger จากระบบไทย, ใช้ AI ทำงานจริง ไม่ใช่ chat"),
        ("ทำไมเรา", "มี LMS เป็น distribution + ลูกค้า corporate trust แล้ว + ทีมเข้าใจทั้ง AI และ business ไทย"),
        ("ขอจาก board/investor", "2-3 hires (Implementation, Sales, Integration Eng) + runway 6 เดือน → ARR ฿7-10M"),
    ]
    y = Inches(2.2)
    for label, text in body:
        add_text(s, Inches(0.6), y, Inches(2.4), Inches(0.4),
                 label, size=13, color=MINT, bold=True, font=FONT_THAI)
        add_text(s, Inches(3.2), y, Inches(9.5), Inches(0.7),
                 text, size=14, color=TEXT_PRIMARY, font=FONT_THAI)
        y += Inches(0.75)
    add_footer(s, page, total)


def slide_problem_intro(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Problem", "ทำไม AI training ในองค์กรไทยถึงล้มเหลว", accent=RED)
    add_text(s, Inches(0.6), Inches(2.3), Inches(12), Inches(1.5),
             '"พนักงานเรียน AI จบแล้ว 100 คน\nแต่ใช้จริงในงานประจำแค่ 8 คน"',
             size=40, color=TEXT_PRIMARY, bold=True, font=FONT_THAI)
    add_text(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.4),
             "— ผลสำรวจ Antiparallel Customer Success, Q1 2026 (n=12 บริษัท)",
             size=12, color=TEXT_SECONDARY, font=FONT_THAI)
    add_divider(s, Inches(0.6), Inches(5.4), Inches(12.1))
    add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.5),
             "ไม่ใช่เพราะคนไม่เก่ง — เพราะระบบในบริษัทไม่รองรับ",
             size=18, color=MINT, font=FONT_THAI)
    add_footer(s, page, total)


def slide_problem_breakdown(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Problem", "3 จุดที่ AI ในองค์กรไทยพัง", accent=RED)
    cols = [
        ("01", "Tools ไม่เชื่อมระบบจริง",
         "ChatGPT/Claude ไม่รู้จัก Line OA, FlowAccount, PromptPay\nพนักงานต้องคัดลอกข้อมูลด้วยมือ ทุกครั้ง ทุกวัน",
         "→ AI กลายเป็น toy ไม่ใช่ tool"),
        ("02", "ความรู้ห้องเรียน ≠ ทักษะหน้างาน",
         "เรียน prompt engineering สวยงาม แต่ทำงานจริงต้องเชื่อม\n5 ระบบ + ส่ง notification + เก็บ audit log",
         "→ Gap ระหว่าง 'รู้' กับ 'ทำได้' กว้างมาก"),
        ("03", "ไม่มี ownership ระดับองค์กร",
         "ใครเป็นเจ้าของ workflow? ใครดูแลเมื่อพัง?\nใครรับผิดถ้า AI ตอบลูกค้าผิด?",
         "→ ไม่มีใครกล้า deploy ของจริง"),
    ]
    col_w = Inches(4.0)
    row_h = Inches(4.0)
    start_x = Inches(0.6)
    start_y = Inches(2.3)
    gap = Inches(0.15)
    for i, (num, title, body, kicker) in enumerate(cols):
        x = start_x + (col_w + gap) * i
        card = add_rect(s, x, start_y, col_w, row_h, BG_SURFACE)
        add_rect(s, x, start_y, col_w, Inches(0.08), RED)
        add_text(s, x + Inches(0.3), start_y + Inches(0.35), Inches(1), Inches(0.4),
                 num, size=12, color=RED, bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.3), start_y + Inches(0.75), Inches(3.5), Inches(0.7),
                 title, size=17, color=TEXT_PRIMARY, bold=True, font=FONT_THAI)
        add_text(s, x + Inches(0.3), start_y + Inches(1.8), Inches(3.5), Inches(1.5),
                 body, size=12, color=TEXT_SECONDARY, font=FONT_THAI)
        add_text(s, x + Inches(0.3), start_y + Inches(3.3), Inches(3.5), Inches(0.5),
                 kicker, size=12, color=MINT, bold=True, font=FONT_THAI)
    add_footer(s, page, total)


def slide_problem_cost(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Problem", "ต้นทุนที่มองไม่เห็นของ 'AI training ที่ไม่ใช้ต่อ'", accent=RED)
    add_text(s, Inches(0.6), Inches(2.2), Inches(7), Inches(0.5),
             "บริษัทขนาด 200 คน ลงทุน AI ต่อปี:",
             size=16, color=TEXT_SECONDARY, font=FONT_THAI)
    rows = [
        ("ค่าเทรน AI ทั้งบริษัท", "฿800,000", "Antiparallel หรือคู่แข่ง"),
        ("ChatGPT Team / Claude Teams", "฿720,000", "$30 × 200 คน × 12 เดือน"),
        ("เวลาที่เสียจากการ copy-paste มือ", "฿2,400,000", "1 ชม./คน/วัน × 200 คน"),
        ("ความผิดพลาดจากการ key มือ", "฿500,000", "ประมาณการ rework + customer churn"),
        ("รวมต้นทุน AI ต่อปี", "฿4,420,000", "ไม่รวมโอกาสที่หาย"),
    ]
    y = Inches(2.9)
    for i, (label, amount, note) in enumerate(rows):
        is_total = i == len(rows) - 1
        bg = BG_SURFACE_2 if is_total else BG_SURFACE
        add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.55), bg)
        add_text(s, Inches(0.85), y + Inches(0.12), Inches(5.5), Inches(0.4),
                 label, size=13, color=TEXT_PRIMARY,
                 bold=is_total, font=FONT_THAI)
        add_text(s, Inches(6.5), y + Inches(0.12), Inches(2.5), Inches(0.4),
                 amount, size=14, color=RED if is_total else TEXT_PRIMARY,
                 bold=True, font=FONT_HEAD, align=PP_ALIGN.RIGHT)
        add_text(s, Inches(9.3), y + Inches(0.12), Inches(3.3), Inches(0.4),
                 note, size=11, color=TEXT_SECONDARY, font=FONT_THAI)
        y += Inches(0.62)
    add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.5),
             "Flow ตัดต้นทุน copy-paste + รวมค่า subscription = ROI > 5× ในปีแรก",
             size=15, color=MINT, bold=True, font=FONT_THAI)
    add_footer(s, page, total)


def slide_why_now(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Why Now", "3 คลื่นที่บรรจบกันใน 2026")
    items = [
        ("AI Literacy ถึงจุดอิ่มตัว",
         "ปี 2024-2025 บริษัทไทย Top 500 ลงทุนเทรน AI ทั่วประเทศ\nตอนนี้ผู้บริหารถามแล้วว่า 'แล้วใช้จริงยังไง?'"),
        ("Model Cost ลด 90% ใน 18 เดือน",
         "GPT-4o-mini, Claude Haiku, Gemini Flash\nทำให้ automation ใช้ AI ราคาเข้าถึงได้สำหรับ SME ไทย"),
        ("Line API + ระบบไทยเปิดมากขึ้น",
         "FlowAccount, PEAK, K-Cash Connect, SCB API\nเปิดให้ developer เชื่อมง่ายกว่าเดิมมาก"),
    ]
    y = Inches(2.3)
    for title, body in items:
        add_rect(s, Inches(0.6), y, Inches(0.08), Inches(1.4), MINT)
        add_text(s, Inches(0.95), y, Inches(11.5), Inches(0.5),
                 title, size=20, color=TEXT_PRIMARY, bold=True, font=FONT_THAI)
        add_text(s, Inches(0.95), y + Inches(0.6), Inches(11.5), Inches(0.9),
                 body, size=13, color=TEXT_SECONDARY, font=FONT_THAI)
        y += Inches(1.65)
    add_footer(s, page, total)


def slide_solution_hero(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Solution", "Antiparallel Flow — AI Workflow Automation for Thai Business")
    add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(1.6),
             "ไม่ใช่ chat — ไม่ใช่ wrapper —\nคือ background automation ที่ทำงานจริงในระบบไทย",
             size=32, color=TEXT_PRIMARY, bold=True, font=FONT_THAI)
    add_text(s, Inches(0.6), Inches(4.4), Inches(12), Inches(0.6),
             "Trigger จากระบบที่ใช้จริง  →  AI ตัดสินใจ  →  ทำงานต่อในระบบเดิม",
             size=18, color=MINT, font=FONT_THAI)
    # 3 chips
    chips = ["Line OA", "Gmail", "Slip OCR", "FlowAccount", "Google Calendar",
             "Slack", "Google Sheets", "PromptPay", "PEAK", "K-Cash"]
    x = Inches(0.6)
    y = Inches(5.4)
    for c in chips:
        # estimate width
        w = Inches(0.3 + 0.13 * len(c))
        chip = add_rect(s, x, y, w, Inches(0.45), BG_SURFACE_2)
        add_text(s, x + Inches(0.1), y + Inches(0.08), w - Inches(0.2), Inches(0.3),
                 c, size=11, color=MINT, font=FONT_HEAD, align=PP_ALIGN.CENTER)
        x += w + Inches(0.15)
        if x > Inches(12):
            x = Inches(0.6)
            y += Inches(0.6)
    add_footer(s, page, total)


def slide_solution_how(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "How It Works", "Flow = Trigger + AI Decision + Action")
    # 4 step flow diagram
    steps = [
        ("01", "TRIGGER", "ลูกค้าทัก Line OA\nใบเสร็จเข้า email\nสลิปโอนเข้ามา"),
        ("02", "AI DECIDE", "Claude / GPT / Gemini\nเข้าใจ context\nตัดสินใจ"),
        ("03", "ACTION", "สร้าง lead ใน CRM\nupdate FlowAccount\nส่ง Slack alert"),
        ("04", "OBSERVE", "Log + audit\nROI dashboard\nretry on failure"),
    ]
    card_w = Inches(2.85)
    card_h = Inches(3.5)
    gap = Inches(0.25)
    start_x = Inches(0.6)
    y = Inches(2.4)
    for i, (num, label, body) in enumerate(steps):
        x = start_x + (card_w + gap) * i
        card = add_rect(s, x, y, card_w, card_h, BG_SURFACE)
        add_rect(s, x, y, card_w, Inches(0.08), MINT)
        add_text(s, x + Inches(0.3), y + Inches(0.4), Inches(1), Inches(0.4),
                 num, size=12, color=MINT, bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.3), y + Inches(0.8), Inches(2.4), Inches(0.5),
                 label, size=17, color=TEXT_PRIMARY, bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.3), y + Inches(1.5), Inches(2.4), Inches(1.8),
                 body, size=12, color=TEXT_SECONDARY, font=FONT_THAI)
        # arrow
        if i < len(steps) - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                         x + card_w + Inches(0.02),
                                         y + Inches(1.5),
                                         Inches(0.2), Inches(0.3)) if False else None
    add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.4),
             "Model-agnostic by design  ·  เปลี่ยน model ได้ทุกเมื่อ  ·  ลด lock-in กับ provider เดียว",
             size=12, color=TEXT_MUTED, font=FONT_THAI, align=PP_ALIGN.CENTER)
    add_footer(s, page, total)


def slide_use_case(prs, page, total, case_num, title_th, trigger, ai_action, result, roi):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, f"Use Case {case_num:02d}", title_th)

    # 3 columns: trigger → AI → result
    blocks = [
        ("TRIGGER", trigger, RED),
        ("AI DECISION", ai_action, MINT),
        ("RESULT", result, ORANGE),
    ]
    col_w = Inches(4.0)
    gap = Inches(0.15)
    start_x = Inches(0.6)
    y = Inches(2.3)
    h = Inches(3.4)
    for i, (label, body, color) in enumerate(blocks):
        x = start_x + (col_w + gap) * i
        add_rect(s, x, y, col_w, h, BG_SURFACE)
        add_rect(s, x, y, Inches(0.08), h, color)
        add_text(s, x + Inches(0.3), y + Inches(0.3), Inches(3), Inches(0.4),
                 label, size=11, color=color, bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.3), y + Inches(0.85), Inches(3.5), Inches(2.4),
                 body, size=13, color=TEXT_PRIMARY, font=FONT_THAI)

    # ROI bar
    roi_y = Inches(6.1)
    add_rect(s, Inches(0.6), roi_y, Inches(12.1), Inches(0.7), BG_SURFACE_2)
    add_text(s, Inches(0.85), roi_y + Inches(0.18), Inches(2), Inches(0.4),
             "ROI", size=11, color=MINT, bold=True, font=FONT_HEAD)
    add_text(s, Inches(2.0), roi_y + Inches(0.15), Inches(10.5), Inches(0.45),
             roi, size=14, color=TEXT_PRIMARY, bold=True, font=FONT_THAI)
    add_footer(s, page, total)


def slide_market_size(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Market", "ตลาดที่เข้าถึงได้จริงในไทย")
    # 3 concentric: TAM / SAM / SOM
    rings = [
        ("TAM", "Total Addressable", "฿18,000M", "บริษัทไทย 750k ราย × AI/automation spend",
         BG_SURFACE_2),
        ("SAM", "Serviceable Addressable", "฿3,200M", "บริษัทไทย 10+ คน ที่ใช้ Line OA + ระบบบัญชี",
         BG_SURFACE),
        ("SOM", "Serviceable Obtainable\n(3 ปี)", "฿420M", "1,000 ลูกค้า × ARR เฉลี่ย ฿420k", MINT),
    ]
    y = Inches(2.2)
    for label, sub, amt, body, color in rings:
        h = Inches(1.4)
        card = add_rect(s, Inches(0.6), y, Inches(12.1), h, BG_SURFACE)
        # accent strip
        strip_color = color if color != MINT else MINT
        add_rect(s, Inches(0.6), y, Inches(0.15), h, strip_color)
        add_text(s, Inches(0.95), y + Inches(0.25), Inches(2.5), Inches(0.4),
                 label, size=22, color=MINT if color == MINT else TEXT_PRIMARY,
                 bold=True, font=FONT_HEAD)
        add_text(s, Inches(0.95), y + Inches(0.75), Inches(4), Inches(0.4),
                 sub, size=11, color=TEXT_SECONDARY, font=FONT_HEAD)
        add_text(s, Inches(5.5), y + Inches(0.35), Inches(3), Inches(0.6),
                 amt, size=28, color=MINT if color == MINT else TEXT_PRIMARY,
                 bold=True, font=FONT_HEAD)
        add_text(s, Inches(8.8), y + Inches(0.45), Inches(4), Inches(0.7),
                 body, size=11, color=TEXT_SECONDARY, font=FONT_THAI)
        y += h + Inches(0.15)
    add_footer(s, page, total)


def slide_competition(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Competition", "ทำไมเรา beat ของที่มีอยู่")
    headers = ["", "ChatGPT Team", "Zapier / Make", "Glean", "Antiparallel Flow"]
    rows = [
        ("Thai integrations (Line OA, FlowAccount)", "—", "บางส่วน", "—", "✓ ครบ"),
        ("ภาษาไทย + code-switching", "ปานกลาง", "—", "อ่อน", "✓ Native"),
        ("Audit trail สำหรับสรรพากร/PDPA", "—", "—", "ปานกลาง", "✓ Built-in"),
        ("Done-for-you setup", "—", "—", "—", "✓ Yes"),
        ("ผูกกับ AI training program", "—", "—", "—", "✓ LMS-linked"),
        ("ราคาเข้าถึงได้ SME ไทย", "ปานกลาง", "ปานกลาง", "แพง", "✓"),
    ]
    col_widths = [Inches(4.3), Inches(1.85), Inches(1.85), Inches(1.85), Inches(2.25)]
    x_offsets = [Inches(0.6)]
    for w in col_widths[:-1]:
        x_offsets.append(x_offsets[-1] + w)

    y = Inches(2.2)
    # header
    for i, h in enumerate(headers):
        is_us = i == 4
        add_rect(s, x_offsets[i], y, col_widths[i], Inches(0.55),
                 MINT if is_us else BG_SURFACE_2)
        add_text(s, x_offsets[i] + Inches(0.15), y + Inches(0.12),
                 col_widths[i] - Inches(0.2), Inches(0.4),
                 h, size=11, color=BG_DARK if is_us else TEXT_PRIMARY,
                 bold=True, font=FONT_HEAD, align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
    y += Inches(0.6)
    for row in rows:
        for i, cell in enumerate(row):
            is_us = i == 4
            bg = BG_SURFACE_2 if is_us else BG_SURFACE
            add_rect(s, x_offsets[i], y, col_widths[i], Inches(0.55), bg)
            color = MINT if (is_us and cell.startswith("✓")) else (RED if cell == "—" else TEXT_PRIMARY)
            add_text(s, x_offsets[i] + Inches(0.15), y + Inches(0.13),
                     col_widths[i] - Inches(0.2), Inches(0.4),
                     cell, size=11, color=color,
                     bold=is_us, font=FONT_THAI,
                     align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
        y += Inches(0.6)
    add_footer(s, page, total)


def slide_moat(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Moat", "4 ชั้นป้องกันที่สะสมยิ่งใช้ยิ่งแน่น")
    moats = [
        ("Thai Integration Library",
         "แต่ละ integration ใช้เวลา build 2-4 สัปดาห์ — ปีแรกเรามี 25 ตัว\nคู่แข่งใหม่ต้องใช้เวลา 12-18 เดือนเพื่อตาม"),
        ("Workflow Template Library",
         "ทุก case study = template ใหม่ — ยิ่งขายยิ่งมี library ใหญ่\nลูกค้าใหม่ deploy ใน 2 สัปดาห์แทน 2 เดือน"),
        ("LMS Distribution Channel",
         "Antiparallel LMS มีลูกค้า corporate active = warm pipeline\nคู่แข่งต้อง cold outreach"),
        ("Customer Data + Workflow Memory",
         "ยิ่งใช้นาน — workflow ผูกกับ business process ขององค์กร\nย้ายออก = พัง ไม่ใช่แค่ไม่สะดวก"),
    ]
    col_w = Inches(6.0)
    row_h = Inches(2.0)
    gap = Inches(0.15)
    start_x = Inches(0.6)
    start_y = Inches(2.2)
    for i, (title, body) in enumerate(moats):
        col = i % 2
        row = i // 2
        x = start_x + (col_w + gap) * col
        y = start_y + (row_h + gap) * row
        add_rect(s, x, y, col_w, row_h, BG_SURFACE)
        add_rect(s, x, y, Inches(0.08), row_h, MINT)
        add_text(s, x + Inches(0.35), y + Inches(0.3), Inches(0.8), Inches(0.4),
                 f"0{i+1}", size=11, color=MINT, bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.35), y + Inches(0.7), Inches(5.5), Inches(0.5),
                 title, size=16, color=TEXT_PRIMARY, bold=True, font=FONT_THAI)
        add_text(s, x + Inches(0.35), y + Inches(1.2), Inches(5.5), Inches(0.8),
                 body, size=11, color=TEXT_SECONDARY, font=FONT_THAI)
    add_footer(s, page, total)


def slide_gtm_overview(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "GTM", "6-month playbook — Bootstrap to ฿7-10M ARR")
    phases = [
        ("Phase 1",
         "M1-M2",
         "Land",
         "5 dealsฟรี/ใกล้ฟรี",
         "ทีม Antiparallel + ลูกค้า LMS + เพื่อนผู้บริหาร + SME ที่ปรึกษา\nDeliverable: 5 case study + 5 MSA",
         MINT),
        ("Phase 2",
         "M3-M4",
         "Productize",
         "Workflow Packs",
         "3 pack: Sales / Finance / HR — ราคา fix\nDeliverable: ขาย 10 ดีล, ARR ฿2.4M",
         ORANGE),
        ("Phase 3",
         "M5-M6",
         "Scale",
         "3 channels parallel",
         "LMS upsell + content inbound + partner pipeline\nDeliverable: 20+ ดีล, 2 enterprise, ARR ฿7-10M",
         RED),
    ]
    y = Inches(2.2)
    for label, period, title, sub, body, color in phases:
        h = Inches(1.45)
        add_rect(s, Inches(0.6), y, Inches(12.1), h, BG_SURFACE)
        add_rect(s, Inches(0.6), y, Inches(0.15), h, color)
        add_text(s, Inches(0.95), y + Inches(0.2), Inches(2), Inches(0.4),
                 label, size=11, color=color, bold=True, font=FONT_HEAD)
        add_text(s, Inches(0.95), y + Inches(0.55), Inches(2), Inches(0.5),
                 period, size=20, color=TEXT_PRIMARY, bold=True, font=FONT_HEAD)
        add_text(s, Inches(2.8), y + Inches(0.2), Inches(3), Inches(0.4),
                 title.upper(), size=11, color=color, bold=True, font=FONT_HEAD)
        add_text(s, Inches(2.8), y + Inches(0.55), Inches(3), Inches(0.5),
                 sub, size=18, color=TEXT_PRIMARY, bold=True, font=FONT_THAI)
        add_text(s, Inches(6.2), y + Inches(0.25), Inches(6.4), Inches(1.1),
                 body, size=12, color=TEXT_SECONDARY, font=FONT_THAI)
        y += h + Inches(0.15)
    add_footer(s, page, total)


def slide_gtm_channel(prs, page, total, num, name, mechanic, pitch, conv_rate, fit_color=MINT):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, f"Channel {num:02d}", name)

    add_text(s, Inches(0.6), Inches(2.3), Inches(2), Inches(0.4),
             "HOW IT WORKS", size=11, color=fit_color, bold=True, font=FONT_HEAD)
    add_text(s, Inches(0.6), Inches(2.7), Inches(12), Inches(1.4),
             mechanic, size=14, color=TEXT_PRIMARY, font=FONT_THAI)

    add_text(s, Inches(0.6), Inches(4.3), Inches(2), Inches(0.4),
             "PITCH", size=11, color=fit_color, bold=True, font=FONT_HEAD)
    add_rect(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(1.4), BG_SURFACE)
    add_rect(s, Inches(0.6), Inches(4.7), Inches(0.08), Inches(1.4), fit_color)
    add_text(s, Inches(0.85), Inches(4.85), Inches(11.7), Inches(1.2),
             pitch, size=13, color=TEXT_PRIMARY, font=FONT_THAI)

    add_text(s, Inches(0.6), Inches(6.3), Inches(3), Inches(0.4),
             "EXPECTED CONVERSION", size=11, color=fit_color, bold=True, font=FONT_HEAD)
    add_text(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.4),
             conv_rate, size=14, color=TEXT_PRIMARY, bold=True, font=FONT_THAI)
    add_footer(s, page, total)


def slide_pricing(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Pricing", "4 tiers — sweet spot ที่ Business")
    tiers = [
        ("STARTER", "SME / Micro", "฿80k", "฿15k/mo", "1 workflow\nemail support\n5k runs/mo", False),
        ("BUSINESS", "Corporate 200+", "฿250k", "฿35k/mo", "3 workflows\nLine support\n50k runs/mo\n1 custom integration", True),
        ("ENTERPRISE", "1,000+ คน", "฿500k-1.5M", "฿80-200k/mo", "Unlimited workflows\nSLA + dedicated eng\nOn-prem option", False),
        ("CONSULTING", "Transformation partner", "฿1.5-3M", "Project-based", "AI Operating Model\n6-month engagement\nC-suite advisory", False),
    ]
    col_w = Inches(3.0)
    gap = Inches(0.1)
    start_x = Inches(0.6)
    y = Inches(2.2)
    h = Inches(4.5)
    for i, (name, sub, setup, monthly, features, highlight) in enumerate(tiers):
        x = start_x + (col_w + gap) * i
        bg = BG_SURFACE_2 if highlight else BG_SURFACE
        add_rect(s, x, y, col_w, h, bg)
        if highlight:
            add_rect(s, x, y, col_w, Inches(0.08), MINT)
            add_text(s, x, y - Inches(0.4), col_w, Inches(0.35),
                     "★ RECOMMENDED", size=10, color=MINT, bold=True,
                     font=FONT_HEAD, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.25), y + Inches(0.3), col_w - Inches(0.4), Inches(0.4),
                 name, size=14, color=MINT if highlight else TEXT_PRIMARY,
                 bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.25), y + Inches(0.7), col_w - Inches(0.4), Inches(0.3),
                 sub, size=10, color=TEXT_MUTED, font=FONT_THAI)
        add_text(s, x + Inches(0.25), y + Inches(1.2), col_w - Inches(0.4), Inches(0.6),
                 setup, size=22, color=TEXT_PRIMARY, bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.25), y + Inches(1.85), col_w - Inches(0.4), Inches(0.35),
                 "setup", size=10, color=TEXT_MUTED, font=FONT_HEAD)
        add_text(s, x + Inches(0.25), y + Inches(2.3), col_w - Inches(0.4), Inches(0.5),
                 monthly, size=16, color=TEXT_PRIMARY, bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.25), y + Inches(2.8), col_w - Inches(0.4), Inches(0.35),
                 "recurring", size=10, color=TEXT_MUTED, font=FONT_HEAD)
        # divider
        add_rect(s, x + Inches(0.25), y + Inches(3.3), col_w - Inches(0.5), Emu(9525), BORDER)
        add_text(s, x + Inches(0.25), y + Inches(3.45), col_w - Inches(0.4), Inches(1.0),
                 features, size=11, color=TEXT_SECONDARY, font=FONT_THAI)
    add_footer(s, page, total)


def slide_unit_economics(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Unit Economics", "เคส Business tier (sweet spot)")
    # Two columns: revenue side / cost side
    # Revenue
    add_text(s, Inches(0.6), Inches(2.2), Inches(6), Inches(0.4),
             "REVENUE PER CUSTOMER (Year 1)", size=12, color=MINT, bold=True, font=FONT_HEAD)
    rev = [
        ("Setup fee (one-time)", "฿250,000"),
        ("Monthly subscription × 12", "฿420,000"),
        ("Total Year 1 Revenue", "฿670,000"),
    ]
    y = Inches(2.7)
    for label, amt in rev:
        is_total = "Total" in label
        bg = BG_SURFACE_2 if is_total else BG_SURFACE
        add_rect(s, Inches(0.6), y, Inches(5.9), Inches(0.55), bg)
        add_text(s, Inches(0.85), y + Inches(0.13), Inches(3.5), Inches(0.4),
                 label, size=12, color=TEXT_PRIMARY, bold=is_total, font=FONT_THAI)
        add_text(s, Inches(4.0), y + Inches(0.13), Inches(2.3), Inches(0.4),
                 amt, size=13, color=MINT if is_total else TEXT_PRIMARY,
                 bold=True, font=FONT_HEAD, align=PP_ALIGN.RIGHT)
        y += Inches(0.62)

    # Cost
    add_text(s, Inches(6.8), Inches(2.2), Inches(6), Inches(0.4),
             "COST PER CUSTOMER (Year 1)", size=12, color=RED, bold=True, font=FONT_HEAD)
    cost = [
        ("Implementation (15 hrs × ฿2k)", "฿30,000"),
        ("API + infra (avg 50k runs/mo)", "฿72,000"),
        ("Support + account mgmt", "฿48,000"),
        ("Sales cost (CAC allocation)", "฿80,000"),
        ("Total Year 1 Cost", "฿230,000"),
    ]
    y = Inches(2.7)
    for label, amt in cost:
        is_total = "Total" in label
        bg = BG_SURFACE_2 if is_total else BG_SURFACE
        add_rect(s, Inches(6.8), y, Inches(5.9), Inches(0.55), bg)
        add_text(s, Inches(7.05), y + Inches(0.13), Inches(3.5), Inches(0.4),
                 label, size=12, color=TEXT_PRIMARY, bold=is_total, font=FONT_THAI)
        add_text(s, Inches(10.2), y + Inches(0.13), Inches(2.3), Inches(0.4),
                 amt, size=13, color=RED if is_total else TEXT_PRIMARY,
                 bold=True, font=FONT_HEAD, align=PP_ALIGN.RIGHT)
        y += Inches(0.62)

    # Bottom metrics
    metrics = [
        ("Gross Margin Y1", "66%"),
        ("Payback Period", "5 เดือน"),
        ("LTV (3yr)", "~฿1.5M"),
        ("LTV / CAC", "18×"),
    ]
    y = Inches(6.2)
    card_w = Inches(2.95)
    for i, (label, val) in enumerate(metrics):
        x = Inches(0.6) + (card_w + Inches(0.1)) * i
        add_rect(s, x, y, card_w, Inches(0.85), BG_SURFACE)
        add_text(s, x + Inches(0.2), y + Inches(0.12), card_w - Inches(0.3), Inches(0.3),
                 label, size=10, color=TEXT_MUTED, font=FONT_HEAD)
        add_text(s, x + Inches(0.2), y + Inches(0.4), card_w - Inches(0.3), Inches(0.5),
                 val, size=22, color=MINT, bold=True, font=FONT_HEAD)
    add_footer(s, page, total)


def slide_financial_projection(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Projection", "6-month financial trajectory")
    months = ["M1", "M2", "M3", "M4", "M5", "M6"]
    customers = ["0", "5*", "5", "10", "15", "22"]
    new_arr = ["—", "—", "฿0.9M", "฿2.4M", "฿4.5M", "฿7.4M"]
    cum_revenue = ["฿0", "฿0", "฿650k", "฿1.2M", "฿2.1M", "฿3.4M"]
    cash_burn = ["฿450k", "฿500k", "฿520k", "฿540k", "฿580k", "฿620k"]

    # Build table
    headers = ["Metric"] + months
    rows = [
        ["Customers (paying)"] + customers,
        ["ARR run-rate"] + new_arr,
        ["Cum. Revenue (cash)"] + cum_revenue,
        ["Monthly Burn"] + cash_burn,
    ]
    col_w = Inches(1.55)
    label_w = Inches(2.85)
    x_start = Inches(0.6)
    y = Inches(2.3)
    # header
    for i, h in enumerate(headers):
        w = label_w if i == 0 else col_w
        x = x_start if i == 0 else x_start + label_w + col_w * (i - 1)
        add_rect(s, x, y, w, Inches(0.55), BG_SURFACE_2)
        add_text(s, x + Inches(0.15), y + Inches(0.12), w - Inches(0.2), Inches(0.4),
                 h, size=12, color=MINT if i > 0 else TEXT_PRIMARY,
                 bold=True, font=FONT_HEAD,
                 align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
    y += Inches(0.6)
    for row in rows:
        for i, cell in enumerate(row):
            w = label_w if i == 0 else col_w
            x = x_start if i == 0 else x_start + label_w + col_w * (i - 1)
            add_rect(s, x, y, w, Inches(0.6), BG_SURFACE)
            add_text(s, x + Inches(0.15), y + Inches(0.15), w - Inches(0.2), Inches(0.4),
                     cell, size=12, color=TEXT_PRIMARY,
                     bold=(i == 0), font=FONT_HEAD,
                     align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
        y += Inches(0.65)

    add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
             "* 5 case studies (free/near-free) → convert to paid in M3",
             size=11, color=TEXT_MUTED, font=FONT_THAI)
    add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
             "Target end of M6: 22 paying customers · ฿7.4M ARR · cash flow positive in M7",
             size=14, color=MINT, bold=True, font=FONT_THAI)
    add_footer(s, page, total)


def slide_team_plan(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Team Plan", "3 hires across 6 months")
    hires = [
        ("M1", "Implementation Engineer",
         "Build workflow ให้ลูกค้า — bottleneck แรก\nต้องเป็น generalist เก่ง API + LLM + product sense",
         "฿80-120k/mo"),
        ("M3", "Sales / Account Manager",
         "หลังมี case study → handle inbound + manage account expansion\nต้องเข้าใจ B2B SaaS sales + ภาษาผู้บริหารไทย",
         "฿60-100k + commission"),
        ("M5", "Integration Engineer",
         "Build platform จริง (no-code editor, queue, observability)\nScale beyond done-for-you",
         "฿100-150k/mo"),
    ]
    y = Inches(2.2)
    for when, role, body, salary in hires:
        h = Inches(1.4)
        add_rect(s, Inches(0.6), y, Inches(12.1), h, BG_SURFACE)
        # When chip
        add_rect(s, Inches(0.6), y, Inches(1.4), h, BG_SURFACE_2)
        add_text(s, Inches(0.6), y + Inches(0.5), Inches(1.4), Inches(0.4),
                 when, size=24, color=MINT, bold=True, font=FONT_HEAD,
                 align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.2), y + Inches(0.2), Inches(8), Inches(0.4),
                 role, size=18, color=TEXT_PRIMARY, bold=True, font=FONT_HEAD)
        add_text(s, Inches(2.2), y + Inches(0.65), Inches(8), Inches(0.7),
                 body, size=12, color=TEXT_SECONDARY, font=FONT_THAI)
        add_text(s, Inches(10.5), y + Inches(0.55), Inches(2.2), Inches(0.4),
                 salary, size=12, color=ORANGE, bold=True, font=FONT_HEAD,
                 align=PP_ALIGN.RIGHT)
        y += h + Inches(0.2)
    add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.4),
             "Founder allocation: M1-M2 sales+build 100%  →  M3-M4 sales+manage  →  M5-M6 enterprise+roadmap",
             size=12, color=TEXT_MUTED, font=FONT_THAI)
    add_footer(s, page, total)


def slide_kpi(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "KPIs", "Health metrics — ทุกเดือนต้องดู")
    kpis = [
        ("Workflow runs / month",
         "Engagement signal — ถ้านิ่ง = ลูกค้าไม่ใช้จริง",
         "Target: 30% MoM growth"),
        ("% customers with ≥ 2 workflows",
         "Expansion signal — ตัวขับ ARR growth ที่แท้จริง",
         "Target: 60% by M6"),
        ("Logo churn",
         "ถ้า > 5%/เดือนใน 6 เดือนแรก = product ไม่ผูกขาดจริง",
         "Target: < 3%/month"),
        ("Buyer NPS (CFO/CEO)",
         "เค้าจ่ายเงิน — ความพอใจเค้าสำคัญที่สุด",
         "Target: > 50"),
        ("User NPS (พนักงาน)",
         "Daily active ของจริง — ขาดมาเร็วถ้าผู้บริหารดัน",
         "Target: > 30"),
        ("Payback period",
         "ดูได้ตั้งแต่ M4 — ถ้า > 9 เดือน = pricing ผิด",
         "Target: < 6 เดือน"),
    ]
    col_w = Inches(4.0)
    row_h = Inches(2.0)
    gap = Inches(0.15)
    start_x = Inches(0.6)
    start_y = Inches(2.2)
    for i, (name, why, target) in enumerate(kpis):
        col = i % 3
        row = i // 3
        x = start_x + (col_w + gap) * col
        y = start_y + (row_h + gap) * row
        add_rect(s, x, y, col_w, row_h, BG_SURFACE)
        add_text(s, x + Inches(0.3), y + Inches(0.25), Inches(3.5), Inches(0.5),
                 name, size=14, color=TEXT_PRIMARY, bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.3), y + Inches(0.85), Inches(3.5), Inches(0.7),
                 why, size=11, color=TEXT_SECONDARY, font=FONT_THAI)
        add_text(s, x + Inches(0.3), y + Inches(1.55), Inches(3.5), Inches(0.4),
                 target, size=11, color=MINT, bold=True, font=FONT_HEAD)
    add_footer(s, page, total)


def slide_risks(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Risks", "สิ่งที่ทำให้เราพัง — และแผนรับมือ", accent=RED)
    risks = [
        ("Model providers ออก native integration",
         "OpenAI / Anthropic เริ่มทำ Line integration เอง",
         "Probability: ต่ำ-กลาง | Impact: สูง",
         "→ Lock-in ด้วย Thai compliance + customer data + LMS bundle\n→ Move up-stack เร็ว: ทำ vertical workflow (sales OS, finance OS)"),
        ("Implementation bottleneck",
         "ขายได้แต่ deliver ไม่ทัน — ลูกค้า cancel",
         "Probability: สูง | Impact: สูง",
         "→ Hire implementation eng วันแรก\n→ จำกัด pipeline ใน M1-M3 ไม่เกิน 5 deals พร้อมกัน"),
        ("API cost spike",
         "ลูกค้าใช้เยอะ margin หาย",
         "Probability: กลาง | Impact: กลาง",
         "→ Per-run pricing tier บน Starter, included runs บน Business\n→ Multi-model routing: ใช้ Haiku/Flash สำหรับ task ง่าย"),
        ("Sale cycle ยาวกว่าคาด",
         "Corporate procurement 3-6 เดือน",
         "Probability: สูง | Impact: กลาง",
         "→ Land via LMS upsell (warm) → bypass procurement\n→ Pilot program ฟรี 30 วัน เพื่อเลี่ยง PO เริ่มต้น"),
    ]
    col_w = Inches(6.0)
    row_h = Inches(2.2)
    gap = Inches(0.15)
    start_x = Inches(0.6)
    start_y = Inches(2.2)
    for i, (risk, what, prob, mitig) in enumerate(risks):
        col = i % 2
        row = i // 2
        x = start_x + (col_w + gap) * col
        y = start_y + (row_h + gap) * row
        add_rect(s, x, y, col_w, row_h, BG_SURFACE)
        add_rect(s, x, y, Inches(0.08), row_h, RED)
        add_text(s, x + Inches(0.3), y + Inches(0.18), Inches(5.5), Inches(0.4),
                 risk, size=14, color=TEXT_PRIMARY, bold=True, font=FONT_THAI)
        add_text(s, x + Inches(0.3), y + Inches(0.62), Inches(5.5), Inches(0.4),
                 what, size=11, color=TEXT_SECONDARY, font=FONT_THAI)
        add_text(s, x + Inches(0.3), y + Inches(1.0), Inches(5.5), Inches(0.4),
                 prob, size=10, color=RED, bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.3), y + Inches(1.4), Inches(5.5), Inches(0.7),
                 mitig, size=10, color=MINT, font=FONT_THAI)
    add_footer(s, page, total)


def slide_roadmap(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Roadmap", "12-month product roadmap")
    quarters = [
        ("Q1 2026", "FOUNDATION",
         "• 5 case studies\n• 8 integrations live\n• Internal admin console\n• Basic observability"),
        ("Q2 2026", "PRODUCTIZE",
         "• 3 Workflow Packs\n• 15 integrations\n• Self-serve onboarding (Starter)\n• Customer dashboard v1"),
        ("Q3 2026", "PLATFORM",
         "• No-code workflow editor\n• 25 integrations (Thai-focus)\n• Multi-model routing\n• Audit/PDPA module"),
        ("Q4 2026", "EXPAND",
         "• Marketplace (3rd-party templates)\n• Enterprise SSO + on-prem\n• Vertical packs (banking, retail)\n• Partner API"),
    ]
    col_w = Inches(3.0)
    gap = Inches(0.1)
    start_x = Inches(0.6)
    y = Inches(2.2)
    h = Inches(4.4)
    for i, (q, theme, body) in enumerate(quarters):
        x = start_x + (col_w + gap) * i
        add_rect(s, x, y, col_w, h, BG_SURFACE)
        add_rect(s, x, y, col_w, Inches(0.08), MINT)
        add_text(s, x + Inches(0.25), y + Inches(0.3), col_w - Inches(0.4), Inches(0.4),
                 q, size=13, color=MINT, bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.25), y + Inches(0.75), col_w - Inches(0.4), Inches(0.5),
                 theme, size=20, color=TEXT_PRIMARY, bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.25), y + Inches(1.5), col_w - Inches(0.4), Inches(2.8),
                 body, size=12, color=TEXT_SECONDARY, font=FONT_THAI)
    add_footer(s, page, total)


def slide_why_us(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Why Us", "ทำไม Antiparallel ทำได้ ทำไมคู่แข่งทำไม่ได้")
    advs = [
        ("Distribution พร้อมใช้",
         "LMS มีลูกค้า corporate active — ไม่ต้อง cold start\nคู่แข่งใหม่ต้องสร้าง trust กับ corporate ไทยจาก 0"),
        ("Domain expertise",
         "ทีมสอน AI ให้องค์กรไทยมาแล้ว เข้าใจ pain ของ HR, ผู้บริหาร, พนักงาน\nรู้ว่าอะไรขายได้ อะไรขายไม่ได้"),
        ("Thai-first DNA",
         "Antiparallel ไทยตั้งแต่วันแรก — ภาษา, integration, สรรพากร, Line\nไม่ใช่ฝรั่ง localize แบบครึ่งๆ กลางๆ"),
        ("Use cases อยู่ในมือแล้ว",
         "Founder มี use case จริงจากบริษัทตัวเอง + ลูกค้า + เพื่อนผู้บริหาร\nเริ่ม build แล้ว validate ได้ทันที"),
    ]
    col_w = Inches(6.0)
    row_h = Inches(2.0)
    gap = Inches(0.15)
    start_x = Inches(0.6)
    start_y = Inches(2.2)
    for i, (title, body) in enumerate(advs):
        col = i % 2
        row = i // 2
        x = start_x + (col_w + gap) * col
        y = start_y + (row_h + gap) * row
        add_rect(s, x, y, col_w, row_h, BG_SURFACE)
        add_rect(s, x, y, Inches(0.08), row_h, MINT)
        add_text(s, x + Inches(0.35), y + Inches(0.3), Inches(5.5), Inches(0.5),
                 title, size=17, color=TEXT_PRIMARY, bold=True, font=FONT_THAI)
        add_text(s, x + Inches(0.35), y + Inches(0.95), Inches(5.5), Inches(1.0),
                 body, size=12, color=TEXT_SECONDARY, font=FONT_THAI)
    add_footer(s, page, total)


def slide_ask(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "The Ask", "สิ่งที่เราขอจาก board / investor")
    add_text(s, Inches(0.6), Inches(2.2), Inches(12), Inches(0.4),
             "RESOURCES", size=11, color=MINT, bold=True, font=FONT_HEAD)
    rows = [
        ("Budget — 6 months operating", "฿4.5M", "3 hires + API/infra + sales/marketing"),
        ("Founder time allocation", "70%", "ลด LMS day-to-day, focus Flow GTM"),
        ("LMS customer access", "Tier 1+2 list", "Workshop intro กับ HR/L&D head"),
        ("Advisor: enterprise SaaS sales", "1 person", "Sounding board สำหรับ enterprise deal"),
    ]
    y = Inches(2.65)
    for label, amt, note in rows:
        add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.7), BG_SURFACE)
        add_text(s, Inches(0.85), y + Inches(0.2), Inches(5.5), Inches(0.4),
                 label, size=13, color=TEXT_PRIMARY, font=FONT_THAI)
        add_text(s, Inches(6.5), y + Inches(0.18), Inches(2.5), Inches(0.4),
                 amt, size=16, color=MINT, bold=True, font=FONT_HEAD,
                 align=PP_ALIGN.RIGHT)
        add_text(s, Inches(9.3), y + Inches(0.22), Inches(3.3), Inches(0.4),
                 note, size=11, color=TEXT_SECONDARY, font=FONT_THAI)
        y += Inches(0.78)

    # Return
    add_text(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.4),
             "EXPECTED OUTCOME (M6)", size=11, color=ORANGE, bold=True, font=FONT_HEAD)
    metrics = [
        ("฿7-10M", "ARR run-rate"),
        ("22+", "Paying customers"),
        ("2+", "Enterprise deals"),
        ("66%", "Gross margin"),
    ]
    y = Inches(6.4)
    card_w = Inches(2.95)
    for i, (val, lab) in enumerate(metrics):
        x = Inches(0.6) + (card_w + Inches(0.1)) * i
        add_rect(s, x, y, card_w, Inches(0.85), BG_SURFACE_2)
        add_text(s, x + Inches(0.2), y + Inches(0.1), card_w - Inches(0.3), Inches(0.5),
                 val, size=22, color=MINT, bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.2), y + Inches(0.55), card_w - Inches(0.3), Inches(0.3),
                 lab, size=10, color=TEXT_SECONDARY, font=FONT_THAI)
    add_footer(s, page, total)


def slide_closing(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    add_rect(s, 0, 0, Inches(0.15), SLIDE_H, MINT)
    add_text(s, Inches(0.8), Inches(0.8), Inches(6), Inches(0.4),
             "CLOSING", size=11, color=MINT, bold=True, font=FONT_HEAD)
    add_text(s, Inches(0.8), Inches(2.4), Inches(12), Inches(1.0),
             "เราไม่ได้สร้างแค่ tool",
             size=48, color=TEXT_PRIMARY, bold=True, font=FONT_THAI)
    add_text(s, Inches(0.8), Inches(3.5), Inches(12), Inches(1.0),
             "เราสร้างชั้น AI ที่ฝังในธุรกิจไทย",
             size=48, color=MINT, bold=True, font=FONT_THAI)
    add_text(s, Inches(0.8), Inches(4.9), Inches(12), Inches(0.5),
             "จาก LMS → ระบบการสอน  →  สู่ Flow → ระบบการทำงาน",
             size=18, color=TEXT_SECONDARY, font=FONT_THAI)
    add_text(s, Inches(0.8), Inches(5.6), Inches(12), Inches(0.5),
             "องค์กรไทยที่ใช้ Antiparallel จะถอดเราออกไม่ได้ในอีก 5 ปี",
             size=18, color=TEXT_SECONDARY, font=FONT_THAI)
    add_divider(s, Inches(0.8), Inches(6.6), Inches(11.7))
    add_text(s, Inches(0.8), Inches(6.75), Inches(8), Inches(0.3),
             "Antiparallel Flow  ·  Let's build it.",
             size=12, color=TEXT_PRIMARY, bold=True, font=FONT_HEAD)
    add_text(s, Inches(10.5), Inches(6.75), Inches(2.5), Inches(0.3),
             "antiparallel.app", size=12, color=MINT, font=FONT_HEAD,
             align=PP_ALIGN.RIGHT)


def slide_appendix_workflows(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Appendix", "20 workflows that sell in Thailand today")
    workflows = [
        "Line OA → Auto reply + CRM lead capture",
        "Slip โอนเงิน → Verify + update order + ส่งใบกำกับ",
        "Email ใบเสร็จ → FlowAccount entry + tax categorization",
        "Resume PDF → Score + rank + auto-schedule interview",
        "Google Meet → Thai notes + action items + Slack",
        "Line group ลูกค้า → Sentiment alert + escalation",
        "Customer email → Categorize + draft reply + assign",
        "Daily sales → Auto-summary deck + email exec team",
        "Stock alert (Sheet) → Reorder + supplier email",
        "Calendar event → Pre-meeting brief from CRM",
        "Support ticket → Suggest reply from knowledge base",
        "Job posting → Multi-channel post (LinkedIn, JobsDB, Line)",
        "Invoice overdue → Polite Thai reminder sequence",
        "ลงเวลาเข้า-ออก → Payroll prep + anomaly detection",
        "OKR check-in → Auto-collect status from Slack/email",
        "Competitor news → Daily Thai summary to leadership",
        "Customer onboarding → Multi-step Line + email journey",
        "Expense report → OCR + policy check + approval routing",
        "Contract → Risk flag + clause comparison vs template",
        "Weekly all-hands → Auto compile updates from each team",
    ]
    col_w = Inches(6.0)
    gap = Inches(0.1)
    items_per_col = 10
    y_start = Inches(2.2)
    for i, wf in enumerate(workflows):
        col = i // items_per_col
        row = i % items_per_col
        x = Inches(0.6) + (col_w + gap) * col
        y = y_start + Inches(0.45) * row
        add_text(s, x, y, Inches(0.4), Inches(0.35),
                 f"{i+1:02d}", size=11, color=MINT, bold=True, font=FONT_HEAD)
        add_text(s, x + Inches(0.45), y, col_w - Inches(0.5), Inches(0.35),
                 wf, size=11, color=TEXT_PRIMARY, font=FONT_THAI)
    add_footer(s, page, total)


def slide_appendix_workshop(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Appendix", "AI Workflow Discovery Workshop — agenda")
    items = [
        ("0:00", "Welcome + objective",
         "Frame: เราจะมาหา 3 workflow ที่ AI ทำได้ในบริษัทคุณ"),
        ("0:10", "Pain point round-robin",
         "ให้แต่ละ manager เล่างานที่เสียเวลาที่สุด 1 อย่าง"),
        ("0:40", "Live demo: build workflow ใน 15 นาที",
         "เลือก 1 pain ที่ทุกคนอินที่สุด → build บนหน้าจอ → ทุกคนเห็น"),
        ("1:00", "Workshop: คัด 3 workflow ที่ ROI สูงสุด",
         "Frame matrix: effort to build × impact when running"),
        ("1:30", "Pricing + next step",
         "Quick Win Package: build 1 workflow ใน 2 สัปดาห์ ฿80k"),
        ("1:45", "Q&A + close",
         "MoU/LOI หาก decision maker อยู่ในห้อง"),
    ]
    y = Inches(2.2)
    for time, title, body in items:
        h = Inches(0.7)
        add_rect(s, Inches(0.6), y, Inches(12.1), h, BG_SURFACE)
        add_text(s, Inches(0.85), y + Inches(0.2), Inches(1), Inches(0.4),
                 time, size=14, color=MINT, bold=True, font=FONT_HEAD)
        add_text(s, Inches(2.2), y + Inches(0.1), Inches(4), Inches(0.4),
                 title, size=13, color=TEXT_PRIMARY, bold=True, font=FONT_THAI)
        add_text(s, Inches(2.2), y + Inches(0.4), Inches(10), Inches(0.3),
                 body, size=11, color=TEXT_SECONDARY, font=FONT_THAI)
        y += h + Inches(0.08)
    add_footer(s, page, total)


def slide_appendix_objections(prs, page, total):
    s = slide_blank(prs)
    add_bg(s)
    content_header(s, "Appendix", "Top objections + คำตอบ")
    objs = [
        ('"ใช้ ChatGPT/Claude เองก็ได้"',
         "เห็นด้วยครับ — แต่ ChatGPT ไม่เชื่อม Line OA, ไม่เข้า FlowAccount, ไม่ออกใบกำกับภาษี\nเรา = glue layer ระหว่าง AI กับระบบไทยที่คุณใช้อยู่"),
        ('"เราใช้ Zapier/Make อยู่แล้ว"',
         "เยี่ยมเลย — เราเป็น layer บน เพราะ Zapier/Make ไม่มี AI native + Thai integration อ่อน\nเรา handle Line OA + ไทย + AI decision logic ที่ Zapier ทำไม่ได้"),
        ('"แพงไป"',
         "เทียบกับเวลาที่เสียจาก copy-paste มือ — บริษัท 200 คน เสีย ฿2.4M/ปี\nFlow ฿420k/ปี = ROI 5×+ ภายใน 6 เดือน"),
        ('"แล้วถ้า AI ตอบผิด ใครรับผิด?"',
         "Human-in-the-loop เป็น default — AI draft, มนุษย์อนุมัติ\nLog ทุก decision เพื่อ audit + adjust prompt ได้ตลอด"),
        ('"เปลี่ยน model ใหม่ทุก 6 เดือน ตามไม่ทัน"',
         "Flow เป็น model-agnostic — เปลี่ยน GPT-5 → Claude → Gemini ได้ใน config\nลูกค้าได้ feature ใหม่ฟรี ไม่ใช่ภาระ"),
    ]
    y = Inches(2.2)
    for q, a in objs:
        h = Inches(0.9)
        add_rect(s, Inches(0.6), y, Inches(12.1), h, BG_SURFACE)
        add_rect(s, Inches(0.6), y, Inches(0.08), h, ORANGE)
        add_text(s, Inches(0.85), y + Inches(0.15), Inches(11.5), Inches(0.4),
                 q, size=13, color=ORANGE, bold=True, font=FONT_THAI)
        add_text(s, Inches(0.85), y + Inches(0.5), Inches(11.5), Inches(0.7),
                 a, size=11, color=TEXT_PRIMARY, font=FONT_THAI)
        y += h + Inches(0.08)
    add_footer(s, page, total)


# ---------- Main ----------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Plan slides (count first to set total)
    # 1. Title
    # 2. TL;DR
    # 3. Executive Summary
    # --- SECTION 1: Problem
    # 4. Section divider
    # 5. Problem intro
    # 6. Problem breakdown
    # 7. Problem cost
    # 8. Why now
    # --- SECTION 2: Solution
    # 9. Section divider
    # 10. Solution hero
    # 11. How it works
    # 12-14. Use cases × 3
    # --- SECTION 3: Market
    # 15. Section divider
    # 16. Market size
    # 17. Competition
    # 18. Moat
    # 19. Why us
    # --- SECTION 4: GTM
    # 20. Section divider
    # 21. GTM overview
    # 22-24. Channels × 3
    # 25. Pricing
    # --- SECTION 5: Financials
    # 26. Section divider
    # 27. Unit economics
    # 28. Financial projection
    # 29. KPIs
    # --- SECTION 6: Team & Roadmap
    # 30. Section divider
    # 31. Team plan
    # 32. Roadmap
    # 33. Risks
    # --- SECTION 7: Close
    # 34. The ask
    # 35. Closing
    # --- Appendix
    # 36. 20 workflows
    # 37. Workshop agenda
    # 38. Objection handling
    TOTAL = 38
    p = 1

    title_slide(prs, p, TOTAL); p += 1
    slide_tldr(prs, p, TOTAL); p += 1
    slide_executive_summary(prs, p, TOTAL); p += 1

    section_divider(prs, 1, "Problem", "ทำไมเราถึงเลือกทำเรื่องนี้", p, TOTAL, accent=RED); p += 1
    slide_problem_intro(prs, p, TOTAL); p += 1
    slide_problem_breakdown(prs, p, TOTAL); p += 1
    slide_problem_cost(prs, p, TOTAL); p += 1
    slide_why_now(prs, p, TOTAL); p += 1

    section_divider(prs, 2, "Solution", "Antiparallel Flow คืออะไร", p, TOTAL); p += 1
    slide_solution_hero(prs, p, TOTAL); p += 1
    slide_solution_how(prs, p, TOTAL); p += 1

    slide_use_case(prs, p, TOTAL, 1,
                   "Line OA → CRM → Sales Follow-up",
                   "ลูกค้าใหม่ทักมาที่ Line OA ของบริษัท\nนอกเวลาทำการ / วันหยุด / สอบถามสินค้า",
                   "Claude อ่านข้อความ + ดู product catalog\nตอบราคา + เก็บข้อมูลลูกค้า + qualify intent\nสร้าง lead ใน CRM พร้อม priority score",
                   "Lead เข้า HubSpot อัตโนมัติ\nSlack แจ้งเซลส์ที่เกี่ยวข้องทันที\nนัด follow-up call ใน Google Calendar",
                   "ลด response time จาก 4 ชม. เหลือ 30 วินาที  ·  Lead conversion +35%  ·  เซลส์ทำงาน high-value ไม่ตอบ FAQ ซ้ำ"); p += 1

    slide_use_case(prs, p, TOTAL, 2,
                   "ใบเสร็จ → FlowAccount → Monthly Tax Prep",
                   "พนักงานส่งใบเสร็จเข้า email ส่วนกลาง\nหรือถ่ายรูปส่ง Line group บัญชี",
                   "GPT-4o-vision อ่านใบเสร็จไทย-อังกฤษ\nแยกหมวด (อาหาร, เดินทาง, สำนักงาน)\nตรวจสอบ vendor + ภาษีซื้อ",
                   "Entry เข้า FlowAccount พร้อม category\nเตือนพนักงานถ้าเอกสารไม่ครบ\nสรุปยอด + ภาษีซื้อรายเดือนให้ CFO",
                   "ลดเวลาทีมบัญชี 60%  ·  ปิดงบเดือนเร็วขึ้น 8 วัน  ·  ลดภาษีซื้อที่ลืม claim ประมาณ ฿80-150k/ปี"); p += 1

    slide_use_case(prs, p, TOTAL, 3,
                   "Resume → AI Screen → Interview Schedule",
                   "Resume PDF เข้ามาผ่าน email หรือ JobsDB\nสำหรับตำแหน่งที่ JD ชัดเจน",
                   "Claude อ่าน resume + เทียบ JD\nคำนวณ score 0-100 พร้อม reasoning\nเช็ค red flag (gap, job hopping)",
                   "Top 10% → auto-email ขอสัมภาษณ์\nนัดผ่าน Calendar slot ที่ HR ว่าง\nส่ง onboarding brief ก่อนสัมภาษณ์",
                   "HR ใช้เวลา screen เหลือ 5 นาที/คน  ·  Time-to-hire ลด 40%  ·  Hiring manager เห็นแต่ candidate ที่ qualify"); p += 1

    section_divider(prs, 3, "Market", "ตลาด + คู่แข่ง + ทำไมเราชนะ", p, TOTAL); p += 1
    slide_market_size(prs, p, TOTAL); p += 1
    slide_competition(prs, p, TOTAL); p += 1
    slide_moat(prs, p, TOTAL); p += 1
    slide_why_us(prs, p, TOTAL); p += 1

    section_divider(prs, 4, "Go-to-Market", "แผนขาย 6 เดือน", p, TOTAL, accent=ORANGE); p += 1
    slide_gtm_overview(prs, p, TOTAL); p += 1

    slide_gtm_channel(prs, p, TOTAL, 1,
                      "LMS Upsell Funnel — channel ที่ขายง่ายที่สุด",
                      "ส่ง email ถึง HR/L&D head ของลูกค้า LMS เดิม เสนอ Workshop ฟรี 2 ชม.\nWorkshop: ระบุ pain point + live demo workflow → ปิดด้วย Quick Win Package ฿80k\nถ้าผลดี 30 วัน → ขยายเป็น Workflow Pack หรือ Enterprise contract",
                      "'บริษัทคุณซื้อคอร์ส AI ให้พนักงาน 100 คนแล้ว — เรียนเสร็จเค้าเอาไปใช้กี่คน?\nผมขอ workshop 2 ชม. กับทีมคุณ เรามาดูว่า workflow ไหน AI ทำได้ — ฟรี ผมจะ build ตัวอย่างให้ดูเลย'",
                      "Email → Workshop 60%  ·  Workshop → Quick Win 30-40%  ·  Quick Win → Enterprise 60-70%",
                      fit_color=MINT); p += 1

    slide_gtm_channel(prs, p, TOTAL, 2,
                      "Founder-led Content — show, don't tell",
                      "Weekly: 'build workflow X ใน 10 นาที' บน YouTube + TikTok + LinkedIn\nCase study (anonymized) จากลูกค้า + dashboard screenshot\nLinkedIn: ผู้บริหารไทย scroll ทุกวัน → DM inbound",
                      "'ลูกค้าผมประหยัด 40 ชม./สัปดาห์ ด้วย workflow นี้' + GIF demo\n→ ผู้บริหารคนอื่นเห็น → comment 'ผมก็อยากได้' → DM → discovery call",
                      "100 views → 1 DM  ·  10 DM → 3 discovery call  ·  3 call → 1 deal  (founder-led 3 เดือนแรก)",
                      fit_color=ORANGE); p += 1

    slide_gtm_channel(prs, p, TOTAL, 3,
                      "Partner Channel — เก็บไว้ Phase 3+",
                      "Partner 3 ประเภท: สำนักงานบัญชี/ทนาย (มี SME portfolio 50-200) → revshare 20%\nIT consultant / SI ไทย → enterprise relationship → co-sell หรือ white label\nLine OA tool partners (Zaapi, Page365) → outgrow path",
                      "'คุณมีลูกค้า SME 100 ราย — ขาย Flow Sales Pack ให้ลูกค้าคุณ ได้ 20% lifetime\nเราดูแล implementation ให้ คุณดูแล relationship'",
                      "Partner sourced deals: target 30% ของ pipeline ใน Year 2",
                      fit_color=RED); p += 1

    slide_pricing(prs, p, TOTAL); p += 1

    section_divider(prs, 5, "Financials", "Unit economics + projection", p, TOTAL, accent=ORANGE); p += 1
    slide_unit_economics(prs, p, TOTAL); p += 1
    slide_financial_projection(prs, p, TOTAL); p += 1
    slide_kpi(prs, p, TOTAL); p += 1

    section_divider(prs, 6, "Execution", "Team, roadmap, risks", p, TOTAL); p += 1
    slide_team_plan(prs, p, TOTAL); p += 1
    slide_roadmap(prs, p, TOTAL); p += 1
    slide_risks(prs, p, TOTAL); p += 1

    slide_ask(prs, p, TOTAL); p += 1
    slide_closing(prs, p, TOTAL); p += 1

    slide_appendix_workflows(prs, p, TOTAL); p += 1
    slide_appendix_workshop(prs, p, TOTAL); p += 1
    slide_appendix_objections(prs, p, TOTAL); p += 1

    out = "/Users/panapat/Brieflylearn/brieflylearn/docs/pitch/Antiparallel_Flow_Strategy_Deck.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Total slides: {p - 1}")


if __name__ == "__main__":
    build()
